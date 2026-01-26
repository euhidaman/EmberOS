"""
Document creation and reading tools for EmberOS.

Supports comprehensive document handling for various formats including:
- Text: .txt, .md, .html, .xml
- Word Processing: .docx, .odt, .rtf
- PDF: .pdf
- Spreadsheets: .xlsx, .ods, .csv
- Presentations: .pptx, .odp
- Images: .jpg, .png, .tiff (for OCR and analysis)
"""

from __future__ import annotations

import os
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Optional
from datetime import datetime

from emberos.tools.base import (
    BaseTool, ToolResult, ToolManifest, ToolParameter,
    ToolCategory, RiskLevel
)
from emberos.tools.registry import register_tool


@register_tool
class DocumentReaderTool(BaseTool):
    """Universal document reader supporting multiple formats."""

    @property
    def manifest(self) -> ToolManifest:
        return ToolManifest(
            name="documents.read",
            description="Read and extract text from various document formats (PDF, DOCX, ODT, images with OCR, etc.)",
            category=ToolCategory.FILESYSTEM,
            icon="📖",
            parameters=[
                ToolParameter(
                    name="path",
                    type="string",
                    description="Path to the document to read",
                    required=True
                ),
                ToolParameter(
                    name="extract_images",
                    type="bool",
                    description="Extract embedded images from document",
                    required=False,
                    default=False
                ),
                ToolParameter(
                    name="ocr",
                    type="bool",
                    description="Apply OCR to image-based PDFs or image files",
                    required=False,
                    default=True
                ),
                ToolParameter(
                    name="max_length",
                    type="int",
                    description="Maximum text length to extract (chars)",
                    required=False,
                    default=50000
                )
            ],
            permissions=["filesystem:read"],
            risk_level=RiskLevel.LOW
        )

    async def execute(self, params: dict[str, Any]) -> ToolResult:
        path = os.path.expanduser(params["path"])
        extract_images = params.get("extract_images", False)
        ocr = params.get("ocr", True)
        max_length = params.get("max_length", 50000)

        try:
            filepath = Path(path)
            if not filepath.exists():
                return ToolResult(success=False, error=f"File not found: {path}")

            extension = filepath.suffix.lower()

            # Determine document type and extract content
            if extension in ['.txt', '.md', '.html', '.htm', '.xml', '.json', '.csv']:
                content = await self._read_text_file(filepath, max_length)
            elif extension in ['.pdf']:
                content = await self._read_pdf(filepath, ocr, extract_images, max_length)
            elif extension in ['.docx', '.doc']:
                content = await self._read_docx(filepath, max_length)
            elif extension in ['.odt', '.rtf']:
                content = await self._read_odt(filepath, max_length)
            elif extension in ['.xlsx', '.xls', '.ods']:
                content = await self._read_spreadsheet(filepath, max_length)
            elif extension in ['.pptx', '.ppt', '.odp']:
                content = await self._read_presentation(filepath, max_length)
            elif extension in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']:
                content = await self._read_image(filepath, ocr, max_length)
            elif extension in ['.epub']:
                content = await self._read_epub(filepath, max_length)
            else:
                return ToolResult(
                    success=False,
                    error=f"Unsupported document format: {extension}"
                )

            return ToolResult(
                success=True,
                data={
                    "path": str(filepath),
                    "format": extension[1:],
                    "content": content,
                    "length": len(content),
                    "truncated": len(content) >= max_length
                }
            )

        except Exception as e:
            return ToolResult(success=False, error=str(e), error_type=type(e).__name__)

    async def _read_text_file(self, filepath: Path, max_length: int) -> str:
        """Read plain text files."""
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            return f.read(max_length)

    async def _read_pdf(self, filepath: Path, ocr: bool, extract_images: bool, max_length: int) -> str:
        """Read PDF files using pdftotext or OCR."""
        try:
            # Try pdftotext first (from poppler-utils)
            result = subprocess.run(
                ['pdftotext', '-layout', str(filepath), '-'],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout[:max_length]
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass

        # Fallback: Try tesseract OCR if enabled
        if ocr:
            try:
                # Convert PDF to images and OCR
                result = subprocess.run(
                    ['tesseract', str(filepath), 'stdout', '-l', 'eng'],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                if result.returncode == 0:
                    return result.stdout[:max_length]
            except (subprocess.TimeoutExpired, FileNotFoundError):
                pass

        return "Could not extract text from PDF. Install 'poppler-utils' or 'tesseract' for PDF reading."

    async def _read_docx(self, filepath: Path, max_length: int) -> str:
        """Read DOCX files."""
        try:
            # Try python-docx
            import docx
            doc = docx.Document(str(filepath))
            text = '\n'.join([para.text for para in doc.paragraphs])
            return text[:max_length]
        except ImportError:
            # Fallback: Use pandoc
            try:
                result = subprocess.run(
                    ['pandoc', '-f', 'docx', '-t', 'plain', str(filepath)],
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                if result.returncode == 0:
                    return result.stdout[:max_length]
            except (subprocess.TimeoutExpired, FileNotFoundError):
                pass
            return "Could not read DOCX. Install 'python-docx' or 'pandoc' for DOCX support."

    async def _read_odt(self, filepath: Path, max_length: int) -> str:
        """Read ODT/RTF files using pandoc."""
        try:
            result = subprocess.run(
                ['pandoc', '-f', 'odt', '-t', 'plain', str(filepath)],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0:
                return result.stdout[:max_length]
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass
        return "Could not read ODT/RTF. Install 'pandoc' for ODT support."

    async def _read_spreadsheet(self, filepath: Path, max_length: int) -> str:
        """Read spreadsheet files."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")

                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                    text_parts.append(row_text)

                    if len(''.join(text_parts)) >= max_length:
                        break

            return ''.join(text_parts)[:max_length]
        except ImportError:
            return "Could not read spreadsheet. Install 'openpyxl' for Excel/ODS support."

    async def _read_presentation(self, filepath: Path, max_length: int) -> str:
        """Read presentation files."""
        try:
            from pptx import Presentation
            prs = Presentation(str(filepath))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n=== Slide {slide_num} ===\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text + '\n')

                if len(''.join(text_parts)) >= max_length:
                    break

            return ''.join(text_parts)[:max_length]
        except ImportError:
            return "Could not read presentation. Install 'python-pptx' for PowerPoint support."

    async def _read_image(self, filepath: Path, ocr: bool, max_length: int) -> str:
        """Read text from images using OCR."""
        if not ocr:
            return "OCR disabled. Cannot extract text from images."

        try:
            result = subprocess.run(
                ['tesseract', str(filepath), 'stdout', '-l', 'eng'],
                capture_output=True,
                text=True,
                timeout=60
            )
            if result.returncode == 0:
                return result.stdout[:max_length]
        except (subprocess.TimeoutExpired, FileNotFoundError):
            return "OCR failed. Install 'tesseract' for image text extraction."

    async def _read_epub(self, filepath: Path, max_length: int) -> str:
        """Read EPUB files."""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            book = epub.read_epub(str(filepath))
            text_parts = []

            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text_parts.append(soup.get_text())

                    if len(''.join(text_parts)) >= max_length:
                        break

            return ''.join(text_parts)[:max_length]
        except ImportError:
            return "Could not read EPUB. Install 'ebooklib' and 'beautifulsoup4' for EPUB support."


@register_tool
class CreateDocumentTool(BaseTool):
    """Create various types of documents from text description."""

    @property
    def manifest(self) -> ToolManifest:
        return ToolManifest(
            name="documents.create",
            description="Create documents in various formats (PDF, DOCX, TXT, MD, HTML, etc.) from content",
            category=ToolCategory.FILESYSTEM,
            icon="📝",
            parameters=[
                ToolParameter(
                    name="path",
                    type="string",
                    description="Path for the document (extension determines format)",
                    required=True
                ),
                ToolParameter(
                    name="content",
                    type="string",
                    description="Content of the document (supports markdown)",
                    required=True
                ),
                ToolParameter(
                    name="title",
                    type="string",
                    description="Document title",
                    required=False,
                    default=""
                ),
                ToolParameter(
                    name="author",
                    type="string",
                    description="Document author",
                    required=False,
                    default=""
                ),
                ToolParameter(
                    name="template",
                    type="string",
                    description="Template style: 'default', 'formal', 'report', 'letter', 'article'",
                    required=False,
                    default="default",
                    choices=["default", "formal", "report", "letter", "article"]
                )
            ],
            permissions=["filesystem:write"],
            risk_level=RiskLevel.LOW
        )

    async def execute(self, params: dict[str, Any]) -> ToolResult:
        path = os.path.expanduser(params["path"])
        content = params["content"]
        title = params.get("title", "")
        author = params.get("author", "")
        template = params.get("template", "default")

        try:
            filepath = Path(path)
            filepath.parent.mkdir(parents=True, exist_ok=True)

            extension = filepath.suffix.lower()

            # Create document based on extension
            if extension in ['.txt']:
                result = await self._create_txt(filepath, content, title)
            elif extension in ['.md']:
                result = await self._create_markdown(filepath, content, title, author)
            elif extension in ['.html', '.htm']:
                result = await self._create_html(filepath, content, title, author, template)
            elif extension in ['.pdf']:
                result = await self._create_pdf(filepath, content, title, author, template)
            elif extension in ['.docx']:
                result = await self._create_docx(filepath, content, title, author, template)
            elif extension in ['.odt']:
                result = await self._create_odt(filepath, content, title, author, template)
            elif extension in ['.rtf']:
                result = await self._create_rtf(filepath, content, title, author)
            elif extension in ['.tex']:
                result = await self._create_latex(filepath, content, title, author)
            elif extension in ['.xml']:
                result = await self._create_xml(filepath, content, title)
            else:
                return ToolResult(
                    success=False,
                    error=f"Unsupported document format for creation: {extension}"
                )

            return ToolResult(
                success=True,
                data={
                    "path": str(filepath),
                    "format": extension[1:],
                    "size": filepath.stat().st_size,
                    "message": f"Created {extension[1:].upper()} document: {filepath.name}"
                }
            )

        except Exception as e:
            return ToolResult(success=False, error=str(e), error_type=type(e).__name__)

    async def _create_txt(self, filepath: Path, content: str, title: str) -> bool:
        """Create plain text file."""
        with open(filepath, 'w', encoding='utf-8') as f:
            if title:
                f.write(f"{title}\n{'=' * len(title)}\n\n")
            f.write(content)
        return True

    async def _create_markdown(self, filepath: Path, content: str, title: str, author: str) -> bool:
        """Create markdown file."""
        with open(filepath, 'w', encoding='utf-8') as f:
            if title:
                f.write(f"# {title}\n\n")
            if author:
                f.write(f"**Author:** {author}  \n")
                f.write(f"**Date:** {datetime.now().strftime('%Y-%m-%d')}\n\n")
            f.write(content)
        return True

    async def _create_html(self, filepath: Path, content: str, title: str, author: str, template: str) -> bool:
        """Create HTML file."""
        html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title or 'Document'}</title>
    <style>
        body {{
            max-width: 800px;
            margin: 50px auto;
            padding: 20px;
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            line-height: 1.6;
            color: #333;
        }}
        h1 {{ color: #2c3e50; border-bottom: 3px solid #3498db; padding-bottom: 10px; }}
        .meta {{ color: #7f8c8d; font-style: italic; margin-bottom: 20px; }}
        pre {{ background: #f4f4f4; padding: 15px; border-left: 4px solid #3498db; overflow-x: auto; }}
        code {{ background: #f4f4f4; padding: 2px 5px; border-radius: 3px; }}
    </style>
</head>
<body>
    <h1>{title or 'Document'}</h1>
    {f'<div class="meta">By {author} | {datetime.now().strftime("%B %d, %Y")}</div>' if author else ''}
    <div class="content">
        {self._markdown_to_html_simple(content)}
    </div>
</body>
</html>"""

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(html_content)
        return True

    def _markdown_to_html_simple(self, text: str) -> str:
        """Simple markdown to HTML conversion."""
        # Very basic conversion - replace newlines with paragraphs
        paragraphs = text.strip().split('\n\n')
        html_paragraphs = [f'<p>{p.replace(chr(10), "<br>")}</p>' for p in paragraphs if p.strip()]
        return '\n'.join(html_paragraphs)

    async def _create_pdf(self, filepath: Path, content: str, title: str, author: str, template: str) -> bool:
        """Create PDF using pandoc or other tools."""
        try:
            # Create temporary markdown file
            with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as tmp:
                tmp.write(f"---\ntitle: {title or 'Document'}\n")
                if author:
                    tmp.write(f"author: {author}\n")
                tmp.write(f"date: {datetime.now().strftime('%Y-%m-%d')}\n---\n\n")
                tmp.write(content)
                tmp_path = tmp.name

            # Convert to PDF using pandoc
            result = subprocess.run(
                ['pandoc', tmp_path, '-o', str(filepath), '--pdf-engine=xelatex'],
                capture_output=True,
                timeout=60
            )

            os.unlink(tmp_path)

            if result.returncode != 0:
                raise Exception("Pandoc conversion failed")

            return True
        except FileNotFoundError:
            # Fallback: Create simple text-based PDF using reportlab
            try:
                from reportlab.lib.pagesizes import letter
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet

                doc = SimpleDocTemplate(str(filepath), pagesize=letter)
                styles = getSampleStyleSheet()
                story = []

                if title:
                    story.append(Paragraph(title, styles['Title']))
                    story.append(Spacer(1, 12))

                if author:
                    story.append(Paragraph(f"By {author}", styles['Italic']))
                    story.append(Spacer(1, 12))

                for para in content.split('\n\n'):
                    if para.strip():
                        story.append(Paragraph(para.replace('\n', '<br/>'), styles['BodyText']))
                        story.append(Spacer(1, 12))

                doc.build(story)
                return True
            except ImportError:
                raise Exception("PDF creation requires 'pandoc' or 'reportlab'. Install with: pip install reportlab")

    async def _create_docx(self, filepath: Path, content: str, title: str, author: str, template: str) -> bool:
        """Create DOCX file."""
        try:
            import docx
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = docx.Document()

            # Add title
            if title:
                title_para = doc.add_heading(title, 0)
                title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Add metadata
            if author or title:
                doc.core_properties.author = author or ""
                doc.core_properties.title = title or ""
                doc.core_properties.created = datetime.now()

            # Add content
            for para in content.split('\n\n'):
                if para.strip():
                    p = doc.add_paragraph(para.strip())

            doc.save(str(filepath))
            return True
        except ImportError:
            raise Exception("DOCX creation requires 'python-docx'. Install with: pip install python-docx")

    async def _create_odt(self, filepath: Path, content: str, title: str, author: str, template: str) -> bool:
        """Create ODT file using pandoc."""
        try:
            with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as tmp:
                if title:
                    tmp.write(f"# {title}\n\n")
                if author:
                    tmp.write(f"**Author:** {author}\n\n")
                tmp.write(content)
                tmp_path = tmp.name

            result = subprocess.run(
                ['pandoc', tmp_path, '-o', str(filepath)],
                capture_output=True,
                timeout=30
            )

            os.unlink(tmp_path)

            if result.returncode != 0:
                raise Exception("Pandoc conversion failed")

            return True
        except FileNotFoundError:
            raise Exception("ODT creation requires 'pandoc'. Install with: sudo pacman -S pandoc")

    async def _create_rtf(self, filepath: Path, content: str, title: str, author: str) -> bool:
        """Create RTF file."""
        rtf_header = r"{\rtf1\ansi\deff0"
        rtf_content = []

        if title:
            rtf_content.append(r"{\b\fs32 " + title + r"}\par\par")

        if author:
            rtf_content.append(r"{\i Author: " + author + r"}\par\par")

        # Convert newlines to RTF paragraphs
        for para in content.split('\n\n'):
            if para.strip():
                rtf_content.append(para.replace('\n', ' ') + r"\par\par")

        rtf_doc = rtf_header + ''.join(rtf_content) + "}"

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(rtf_doc)
        return True

    async def _create_latex(self, filepath: Path, content: str, title: str, author: str) -> bool:
        """Create LaTeX file."""
        latex_content = r"""\documentclass{article}
\usepackage[utf8]{inputenc}
\usepackage{hyperref}

"""
        if title:
            latex_content += f"\\title{{{title}}}\n"
        if author:
            latex_content += f"\\author{{{author}}}\n"

        latex_content += f"\\date{{{datetime.now().strftime('%B %d, %Y')}}}\n\n"
        latex_content += r"\begin{document}" + "\n\n"

        if title:
            latex_content += r"\maketitle" + "\n\n"

        # Escape special LaTeX characters
        content_escaped = content.replace('\\', '\\textbackslash{}')
        content_escaped = content_escaped.replace('&', '\\&').replace('%', '\\%')
        content_escaped = content_escaped.replace('$', '\\$').replace('#', '\\#')
        content_escaped = content_escaped.replace('_', '\\_').replace('{', '\\{').replace('}', '\\}')

        latex_content += content_escaped + "\n\n"
        latex_content += r"\end{document}"

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(latex_content)
        return True

    async def _create_xml(self, filepath: Path, content: str, title: str) -> bool:
        """Create XML file."""
        xml_content = f"""<?xml version="1.0" encoding="UTF-8"?>
<document>
    <metadata>
        <title>{title or 'Document'}</title>
        <created>{datetime.now().isoformat()}</created>
    </metadata>
    <content>
        <![CDATA[{content}]]>
    </content>
</document>"""

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(xml_content)
        return True


@register_tool
class ConvertDocumentTool(BaseTool):
    """Convert documents between different formats."""

    @property
    def manifest(self) -> ToolManifest:
        return ToolManifest(
            name="documents.convert",
            description="Convert documents between formats (e.g., PDF to DOCX, MD to HTML, etc.)",
            category=ToolCategory.FILESYSTEM,
            icon="🔄",
            parameters=[
                ToolParameter(
                    name="source",
                    type="string",
                    description="Source document path",
                    required=True
                ),
                ToolParameter(
                    name="destination",
                    type="string",
                    description="Destination path (extension determines output format)",
                    required=True
                ),
                ToolParameter(
                    name="options",
                    type="dict",
                    description="Conversion options (depends on format)",
                    required=False,
                    default={}
                )
            ],
            permissions=["filesystem:read", "filesystem:write"],
            risk_level=RiskLevel.LOW
        )

    async def execute(self, params: dict[str, Any]) -> ToolResult:
        source = os.path.expanduser(params["source"])
        destination = os.path.expanduser(params["destination"])
        options = params.get("options", {})

        try:
            src_path = Path(source)
            dst_path = Path(destination)

            if not src_path.exists():
                return ToolResult(success=False, error=f"Source file not found: {source}")

            dst_path.parent.mkdir(parents=True, exist_ok=True)

            # Try pandoc for conversion
            try:
                result = subprocess.run(
                    ['pandoc', str(src_path), '-o', str(dst_path)],
                    capture_output=True,
                    text=True,
                    timeout=60
                )

                if result.returncode == 0:
                    return ToolResult(
                        success=True,
                        data={
                            "source": str(src_path),
                            "destination": str(dst_path),
                            "size": dst_path.stat().st_size,
                            "message": f"Converted {src_path.suffix} to {dst_path.suffix}"
                        }
                    )
                else:
                    return ToolResult(
                        success=False,
                        error=f"Conversion failed: {result.stderr}"
                    )
            except FileNotFoundError:
                return ToolResult(
                    success=False,
                    error="Document conversion requires 'pandoc'. Install with: sudo pacman -S pandoc"
                )

        except Exception as e:
            return ToolResult(success=False, error=str(e), error_type=type(e).__name__)

